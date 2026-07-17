import base64
import json
import logging
import mimetypes
import re
import urllib.request
from pathlib import Path

from app.core.config import settings

TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".html", ".htm"}
PDF_EXTENSIONS = {".pdf"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
OFFICE_EXTENSIONS = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"}

logger = logging.getLogger("app")


class ChunkingClient:
    """多模态解析与切分 API 的预留位置。

    当前：
      - 文本文件直接读取内容
      - PDF 使用 PyMuPDF 提取文本；扫描页可选走 DashScope OCR/VL
      - 图片可选走 DashScope OCR/VL
      - docx/xlsx/pptx 使用本地解析库提取正文
    """

    def detect_modality(self, filename: str, mime_type: str = "") -> str:
        suffix = Path(filename or "").suffix.lower()
        mime_type = (mime_type or "").lower()
        if suffix in PDF_EXTENSIONS or mime_type == "application/pdf":
            return "pdf"
        if suffix in IMAGE_EXTENSIONS or mime_type.startswith("image/"):
            return "image"
        if suffix in TEXT_EXTENSIONS or mime_type.startswith("text/"):
            return "text"
        if suffix in OFFICE_EXTENSIONS:
            return "office"
        return "text"

    def extract_placeholder(self, file_path: Path, mime_type: str, description: str = "") -> dict:
        modality = self.detect_modality(file_path.name, mime_type)
        description = (description or "").strip()
        if modality == "text":
            content, status = self._extract_text_file(file_path)
        elif modality == "pdf":
            content, status = self._extract_pdf(file_path)
        elif modality == "image":
            content, status = self._extract_image(file_path, mime_type)
        elif modality == "office":
            content, status = self._extract_office(file_path)
        else:
            content = description or f"文件 {file_path.name} 已上传，等待多模态解析。"
            status = "placeholder"
        content = self._merge_description(description, content)
        content = self._clip(content)
        return {"modality": modality, "content": content, "extraction_status": status}

    # ── LangChain 切分参数 ────────────────────────────────────────
    # chunk_size：每块目标字符数。中文 900 字 ≈ 333 tokens，
    #   加上标题前缀后仍远低于 DashScope text-embedding-v3 的 8192 token 上限。
    # chunk_overlap：相邻块重叠字符数，保留跨块上下文。
    CHUNK_SIZE = 900
    CHUNK_OVERLAP = 120

    # RecursiveCharacterTextSplitter 按优先级逐级尝试以下分隔符：
    #   双换行（段落）→ 单换行 → 中文句末标点 → 英文句末标点 → 空格 → 单字符
    _SEPARATORS = ["\n\n", "\n", "。", "！", "？", "；", ".", "!", "?", ";", " ", ""]
    _PDF_PAGE_PATTERN = re.compile(r"【第\s*(\d+)\s*页】")

    def chunk_document(
        self,
        source_id: str,
        title: str,
        content: str,
        modality: str,
        section_path: str = "",
        file_path: str = "",
        mime_type: str = "",
    ) -> list[dict]:
        """使用 LangChain RecursiveCharacterTextSplitter 切分文档。

        每个 chunk 作为向量数据库中的独立记录，支持 PDF / Word / 纯文本等任意格式。
        切分策略：优先在段落/句子边界断开，保证语义完整性。
        """
        try:
            from langchain_text_splitters import RecursiveCharacterTextSplitter
        except ImportError:
            logger.warning(
                "langchain-text-splitters 未安装，降级为固定长度切分。"
                "请运行 pip install langchain-text-splitters==1.1.2"
            )
            return self._chunk_fixed(source_id, title, content, modality, section_path)

        text = (content or "").strip()
        if not text:
            text = f"{title} 已上传，等待解析。"

        if modality == "image":
            return self._chunk_image_asset(
                source_id=source_id,
                title=title,
                text=text,
                section_path=section_path,
                file_path=file_path,
                mime_type=mime_type,
            )

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.CHUNK_SIZE,
            chunk_overlap=self.CHUNK_OVERLAP,
            separators=self._SEPARATORS,
            length_function=len,
            is_separator_regex=False,
        )
        if modality == "pdf":
            page_assets = self._render_pdf_page_assets(source_id, file_path)
            page_chunks = self._chunk_pdf_pages(
                splitter,
                source_id=source_id,
                title=title,
                text=text,
                modality=modality,
                section_path=section_path,
                page_assets=page_assets,
            )
            if page_chunks:
                return page_chunks

        raw_chunks = splitter.split_text(text)

        chunks = []
        for index, chunk_text in enumerate(raw_chunks):
            chunk_text = chunk_text.strip()
            if not chunk_text:
                continue
            chunks.append(
                {
                    "chunk_id": f"{source_id}-chunk-{index:04d}",
                    "source_id": source_id,
                    "chunk_index": index,
                    "section_path": section_path or title,
                    "content": chunk_text,
                    "modality": modality,
                    "parser": "langchain_recursive_splitter",
                }
            )
        return chunks

    def _chunk_image_asset(
        self,
        *,
        source_id: str,
        title: str,
        text: str,
        section_path: str,
        file_path: str,
        mime_type: str,
    ) -> list[dict]:
        return [
            {
                "chunk_id": f"{source_id}-chunk-0000",
                "source_id": source_id,
                "chunk_index": 0,
                "section_path": section_path or title,
                "content": text,
                "modality": "image",
                "parser": "image_multimodal_asset",
                "page_number": 0,
                "asset_path": file_path,
                "asset_mime_type": mime_type or "image/png",
                "embedding_modality": "image" if file_path else "text",
            }
        ]

    def _render_pdf_page_assets(self, source_id: str, file_path: str) -> dict[str, str]:
        if not file_path:
            return {}
        path = Path(file_path)
        if not path.exists():
            return {}
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {}

        output_dir = settings.asset_dir / source_id
        output_dir.mkdir(parents=True, exist_ok=True)
        page_assets: dict[str, str] = {}
        try:
            with fitz.open(path) as doc:
                for page_index, page in enumerate(doc, start=1):
                    asset_path = output_dir / f"page-{page_index:04d}.jpg"
                    if not asset_path.exists():
                        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                        pix.save(str(asset_path))
                    page_assets[str(page_index)] = str(asset_path)
        except Exception as exc:
            logger.warning("PDF page asset rendering failed for %s: %s", file_path, exc)
            return {}
        return page_assets

    def _chunk_pdf_pages(
        self,
        splitter,
        *,
        source_id: str,
        title: str,
        text: str,
        modality: str,
        section_path: str,
        page_assets: dict[str, str] | None = None,
    ) -> list[dict]:
        pages = self._split_pdf_pages(text)
        if not pages:
            return []

        chunks = []
        index = 0
        base_section = section_path or title
        page_assets = page_assets or {}
        for page_no, page_text in pages:
            page_asset = page_assets.get(str(page_no), "")
            page_section = f"{base_section} / 第 {page_no} 页"
            for raw_chunk in splitter.split_text(page_text):
                chunk_text = raw_chunk.strip()
                if not chunk_text:
                    continue
                chunks.append(
                    {
                        "chunk_id": f"{source_id}-chunk-{index:04d}",
                        "source_id": source_id,
                        "chunk_index": index,
                        "section_path": page_section,
                        "content": chunk_text,
                        "modality": modality,
                        "parser": "langchain_pdf_page_splitter",
                        "page_number": int(page_no),
                        "asset_path": page_asset,
                        "asset_mime_type": "image/jpeg" if page_asset else "",
                        "embedding_modality": "pdf_page" if page_asset else "text",
                    }
                )
                index += 1
        return chunks

    def _split_pdf_pages(self, text: str) -> list[tuple[str, str]]:
        matches = list(self._PDF_PAGE_PATTERN.finditer(text))
        if not matches:
            return []

        pages = []
        for index, match in enumerate(matches):
            start = match.end()
            end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
            page_text = text[start:end].strip()
            if page_text:
                pages.append((match.group(1), page_text))
        return pages

    def _chunk_fixed(
        self,
        source_id: str,
        title: str,
        content: str,
        modality: str,
        section_path: str = "",
    ) -> list[dict]:
        """固定长度切分（langchain-text-splitters 未安装时的降级实现）。"""
        text = (content or "").strip() or f"{title} 已上传，等待解析。"
        max_chars, overlap = 420, 60
        chunks, start, index = [], 0, 0
        while start < len(text):
            end = min(len(text), start + max_chars)
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(
                    {
                        "chunk_id": f"{source_id}-chunk-{index:04d}",
                        "source_id": source_id,
                        "chunk_index": index,
                        "section_path": section_path or title,
                        "content": chunk_text,
                        "modality": modality,
                        "parser": "fixed_chunker_fallback",
                    }
                )
                index += 1
            if end >= len(text):
                break
            start = max(0, end - overlap)
        return chunks

    # ── 解析实现 ──────────────────────────────────────────────────

    def _extract_text_file(self, file_path: Path) -> tuple[str, str]:
        try:
            return file_path.read_text(encoding="utf-8"), "parsed"
        except UnicodeDecodeError:
            return file_path.read_text(encoding="gbk", errors="replace"), "parsed"
        except Exception as exc:
            logger.warning("Text extraction failed for %s: %s", file_path, exc)
            return f"文本文件 {file_path.name} 已上传，但解析失败：{exc}", "placeholder"

    def _extract_pdf(self, file_path: Path) -> tuple[str, str]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return (
                f"PDF 文件 {file_path.name} 已上传，但未安装 PyMuPDF，无法提取正文。",
                "parser_dependency_missing",
            )

        parts: list[str] = []
        used_ocr = False
        try:
            with fitz.open(file_path) as doc:
                for page_index, page in enumerate(doc, start=1):
                    text = self._extract_pdf_page_text(page)
                    if not text and self._ocr_enabled():
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                        text = self._ocr_bytes(
                            pix.tobytes("png"),
                            "image/png",
                            f"请提取这页 PDF 图片中的全部文字，并保持自然段顺序。页码：{page_index}",
                        ).strip()
                        used_ocr = used_ocr or bool(text)
                    if text:
                        parts.append(f"【第 {page_index} 页】\n{text}")
        except Exception as exc:
            logger.warning("PDF extraction failed for %s: %s", file_path, exc)
            return f"PDF 文件 {file_path.name} 已上传，但解析失败：{exc}", "placeholder"

        if parts:
            return "\n\n".join(parts), "parsed_with_ocr" if used_ocr else "parsed"
        return (
            f"PDF 文件 {file_path.name} 已上传，但未提取到文本。若这是扫描版 PDF，请配置 OCR 后重新入库。",
            "pending_pdf_ocr",
        )

    def _extract_pdf_page_text(self, page) -> str:
        blocks = []
        try:
            for block in page.get_text("blocks"):
                if len(block) < 5:
                    continue
                x0, y0, _x1, _y1, text = block[:5]
                text = str(text).strip()
                if text:
                    blocks.append((round(float(y0), 1), round(float(x0), 1), text))
        except Exception:
            blocks = []

        if blocks:
            blocks.sort(key=lambda item: (item[0], item[1]))
            return "\n".join(item[2] for item in blocks).strip()
        return page.get_text("text").strip()

    def _extract_image(self, file_path: Path, mime_type: str) -> tuple[str, str]:
        if not self._ocr_enabled():
            return (
                f"图片文件 {file_path.name} 已上传。当前未启用 OCR，无法自动提取图片文字。",
                "pending_vision_parser",
            )
        try:
            actual_mime = mime_type or mimetypes.guess_type(file_path.name)[0] or "image/png"
            text = self._ocr_bytes(
                file_path.read_bytes(),
                actual_mime,
                "请提取图片中的全部文字；如果没有可见文字，请描述图片中的关键对象、场景和信息。",
            ).strip()
            if text:
                return text, "parsed_with_ocr"
        except Exception as exc:
            logger.warning("Image OCR failed for %s: %s", file_path, exc)
            return f"图片文件 {file_path.name} 已上传，但 OCR 失败：{exc}", "placeholder"
        return f"图片文件 {file_path.name} 已上传，但 OCR 未返回可用文本。", "placeholder"

    def _extract_office(self, file_path: Path) -> tuple[str, str]:
        suffix = file_path.suffix.lower()
        try:
            if suffix == ".docx":
                return self._extract_docx(file_path), "parsed"
            if suffix == ".xlsx":
                return self._extract_xlsx(file_path), "parsed"
            if suffix == ".pptx":
                return self._extract_pptx(file_path), "parsed"
            return (
                f"Office 文件 {file_path.name} 已上传。当前仅支持 docx/xlsx/pptx；旧版 doc/xls/ppt 请先另存为新版格式。",
                "unsupported_office_format",
            )
        except ImportError as exc:
            return (
                f"Office 文件 {file_path.name} 已上传，但缺少解析依赖：{exc}",
                "parser_dependency_missing",
            )
        except Exception as exc:
            logger.warning("Office extraction failed for %s: %s", file_path, exc)
            return f"Office 文件 {file_path.name} 已上传，但解析失败：{exc}", "placeholder"

    def _extract_docx(self, file_path: Path) -> str:
        from docx import Document

        doc = Document(file_path)
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table_index, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"【表格 {table_index}】\n" + "\n".join(rows))
        return "\n\n".join(parts) or f"DOCX 文件 {file_path.name} 未提取到文本。"

    def _extract_xlsx(self, file_path: Path) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        try:
            parts = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value not in (None, "")]
                    if values:
                        rows.append(" | ".join(values))
                if rows:
                    parts.append(f"【工作表：{ws.title}】\n" + "\n".join(rows))
            return "\n\n".join(parts) or f"XLSX 文件 {file_path.name} 未提取到文本。"
        finally:
            wb.close()

    def _extract_pptx(self, file_path: Path) -> str:
        from pptx import Presentation

        prs = Presentation(file_path)
        parts = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            if texts:
                parts.append(f"【第 {slide_index} 页】\n" + "\n".join(texts))
        return "\n\n".join(parts) or f"PPTX 文件 {file_path.name} 未提取到文本。"

    # ── OCR / 通用工具 ────────────────────────────────────────────

    def _ocr_enabled(self) -> bool:
        return (
            settings.parser_ocr_provider.lower() == "dashscope"
            and bool(settings.parser_ocr_api_key or settings.embedding_api_key)
        )

    def _ocr_bytes(self, image_bytes: bytes, mime_type: str, prompt: str) -> str:
        api_key = settings.parser_ocr_api_key or settings.embedding_api_key
        chat_url = f"{settings.parser_ocr_base_url.rstrip('/')}/chat/completions"
        image_b64 = base64.b64encode(image_bytes).decode("ascii")
        payload = json.dumps(
            {
                "model": settings.parser_ocr_model,
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime_type};base64,{image_b64}"},
                            },
                            {"type": "text", "text": prompt},
                        ],
                    }
                ],
            },
            ensure_ascii=False,
        ).encode("utf-8")
        req = urllib.request.Request(
            chat_url,
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=settings.parser_ocr_timeout) as resp:
            result = json.loads(resp.read().decode("utf-8"))
        return result["choices"][0]["message"]["content"]

    def _merge_description(self, description: str, content: str) -> str:
        if not description:
            return content
        if not content or content == description:
            return description
        return f"用户说明：\n{description}\n\n解析正文：\n{content}"

    def _clip(self, content: str) -> str:
        limit = max(1000, int(settings.parser_max_chars or 200_000))
        if len(content) <= limit:
            return content
        return content[:limit] + f"\n\n【内容过长，已截断到 {limit} 字符】"


chunking_client = ChunkingClient()
